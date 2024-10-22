import streamlit as st
from openai import OpenAI
import openai
import os
import json
from pptx import Presentation
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import re
import unicodedata


# Funktion zum Initialisieren der session_state Variablen
def initialize_session_state():
    if 'new_title' not in st.session_state:
        st.session_state.new_title = ""
    if 'ai_api_info' not in st.session_state:
        st.session_state.ai_api_info = ""
    if 'new_header' not in st.session_state:
        st.session_state.new_header = ""
    if 'new_content_focus' not in st.session_state:
        st.session_state.new_content_focus = ""
    if 'new_doctype' not in st.session_state:
        st.session_state.new_doctype = "Sales Folien"
    if 'new_chapter_count' not in st.session_state:
        st.session_state.new_chapter_count = 8

    if 'new_context' not in st.session_state:
        st.session_state.new_context = ""
    if 'new_stakeholder' not in st.session_state:
        st.session_state.new_stakeholder = "Technisches Fachpersonal"

    if 'toc_list' not in st.session_state:
        st.session_state.toc_list = []
    if 'kapitel_header' not in st.session_state:
        st.session_state.kapitel_header = []
    if 'kapitel_info' not in st.session_state:
        st.session_state.kapitel_info = []
    if 'kapitel_inhalt' not in st.session_state:
        st.session_state.kapitel_inhalt = []
    if 'kapitel_prompt' not in st.session_state:
        st.session_state.kapitel_prompt = []
    if 'image_prompt' not in st.session_state:
        st.session_state.kapitel_prompt = []
    if "openai_model" not in st.session_state:
        st.session_state["openai_model"] = ""

# Initialisiere session_state
initialize_session_state()

# Funktionen ---------------------------------------------------------------------------------------------------------------------------
#hole dir den ai_key entweder aus der OS Umgebungsvariable oder dem Streamlit Secret Vault
#Azure OpenAI Connection
def get_oai_client():
    if "AZURE_OPENAI_API_KEY" in os.environ:
        client = openai.AzureOpenAI(
            api_key=os.environ["AZURE_OPENAI_API_KEY"],
            api_version="2023-03-15-preview",
            azure_endpoint="https://mlu-azure-openai-service-sw.openai.azure.com/"
        )
        st.session_state["openai_model"] = "gpt-4o-mini-sw"   
    elif "OPENAI_API_KEY" in st.secrets:
        client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
        st.session_state["openai_model"] = "gpt-4o-mini"
        #st.session_state.ai_api_info="powered by OpenAI"
    else:
        raise ValueError("Kein gültiger API-Schlüssel gefunden.")
    return client


client = get_oai_client()

def generate_toc(new_title, new_doctype, new_content_focus, new_chapter_count):

    response = client.chat.completions.create(
        model=st.session_state["openai_model"],
        messages=[
            {"role": "system", "content": "Du bist ein Assistent, der eine Story für eine Powerpointpräsentation erstellt"},
            {"role":"user" , "content": "Erstelle Inhalte für Powerpointfolien mit Folientitel und Stichpunkten für den Folieninhalt für eine Präsentation vom Typ " + new_doctype + " zum Thema " + new_title + " mit etwa " + str(new_chapter_count) + " Folien."},
            {"role":"user" , "content": "Der inhaltliche Schwerpunkt sollte auf folgende Punkte gesetzt werden: " + new_content_focus},
            {"role": "user", "content": "Erstelle zu jeder Folie Notizen, welche kurz erklären, was zu dieser Folie auf der Tonspur erzählt werden sollte"},
            {"role": "user", "content": "Erstelle zu jeder Folie einen Chatprompt, um den Folieninhalt ggf. noch mal neu zu generieren"},
            {"role": "user", "content": "Erstelle zu jeder Folie einen Chatprompt, um ein Bild für diese Inhaltsfolie zu generieren"}
        ],
        functions=[
            {
                "name": "generate_toc",
                "description": "Generates a table of contents with notes for a given topic",
                "parameters": {
                    "type": "object",                    
                    "properties": {
                        "toc": {
                            "type": "array",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {
                                        "type": "string",
                                        "description": "The title of the slide"
                                    },
                                    "content_text": {
                                        "type": "string",
                                        "description": "Up to 8 Content bulletpoints for this slide"
                                    },
                                    "help_text": {
                                        "type": "string",
                                        "description": "Notes for the slide"
                                    },
                                    "prompt_text": {
                                        "type": "string",
                                        "description": "A prompt to make a chatbot like chatgpt generate up to 6 bulletpoints as content for this specific slide."
                                    },
                                    "image_prompt_text": {
                                        "type": "string",
                                        "description": "A prompt to make a chatbot like chatgpt generate an image for this specific slide."
                                    }
                                },
                                "required": ["title", "content_text", "help_text", "prompt_text", "image_prompt_text"]
                            },
                            "description": "The generated table of contents with help texts"
                        }
                    },
                    "required": ["toc"]
                }
            }
        ],
        function_call="auto"
    )

    # Parse the response
    toc = json.loads(response.choices[0].message.function_call.arguments)
    toc_list = toc["toc"]

    st.write("Verwendete AI Tokens zur Erstellung der Folien: " + str(response.usage.total_tokens))
    #### debug
    st.write(toc_list)

    return toc_list

# Haupseite -------------------------------------------------------------------------------------------------------------------------------------
st.set_page_config(page_title="PowerPoint AI", page_icon=":mechanical_arm:", layout="wide")
st.title("PowerPoint AI")

# LLM Info abrufen
if "AZURE_OPENAI_API_KEY" in os.environ:
    st.session_state.ai_api_info="Azure OpenAI - Region Europa"
elif "OPENAI_API_KEY" in st.secrets:
    st.session_state.ai_api_info="powered by OpenAI"
else:
    st.session_state.ai_api_info="OpenAI"

st.write(st.session_state.ai_api_info)



# Funktionen für die PPT Verarbeitung ---------------------------------------------------------------
def open_pptx_template():
    # Pfad zur PowerPoint-Datei
    pptx_path = "msg_digital_template.pptx"
    
    try:
        # Öffnen der PowerPoint-Präsentation
        presentation = Presentation(pptx_path)
        # Hier können Sie weitere Aktionen mit der geöffneten Präsentation durchführen
        
        return presentation
    except Exception as e:
        st.write(f"Fehler beim Öffnen der Präsentation: {e}")
        return None


def generate_ppt(presentation, title, subtitle):
    
      
    # Ändern des Titels und Untertitels
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text == "*title*":
                    shape.text = title
                elif shape.text == "*subtitle*":
                    shape.text = subtitle
    
    # Speichern der Präsentation in einem BytesIO-Objekt
    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    pptx_file = buffer

    # Bereitstellen der Datei zum Download
    st.sidebar.download_button(
        label="PowerPoint-Präsentation herunterladen",
        data=pptx_file,
        file_name="Präsentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

aws_migration_slides = [
    {
        "titel": "Einführung in AWS Migration",
        "inhalt": [
            "Definition von Cloud-Migration",
            "Vorteile der Migration zu AWS",
            "Überblick über den AWS-Migrationsansatz",
            "Wichtige Überlegungen vor der Migration",
            "AWS Migration Competency Partner"
        ]
    },
    {
        "titel": "AWS Migration-Strategien",
        "inhalt": [
            "Rehosting (Lift-and-Shift)",
            "Replatforming",
            "Refactoring / Re-architecting",
            "Repurchasing",
            "Retire und Retain Strategien"
        ]
    },
    {
        "titel": "AWS Migration Tools",
        "inhalt": [
            "AWS Migration Hub",
            "AWS Application Discovery Service",
            "AWS Database Migration Service (DMS)",
            "AWS Server Migration Service (SMS)",
            "AWS CloudEndure Migration"
        ]
    },
    {
        "titel": "Phasen der AWS Migration",
        "inhalt": [
            "Assess: Bewertung der Migrationsbereitschaft",
            "Mobilize: Vorbereitung der Organisation",
            "Migrate & Modernize: Durchführung der Migration",
            "Operate & Optimize: Betrieb in der Cloud",
            "Sicherheit und Compliance während der Migration"
        ]
    },
    {
        "titel": "Best Practices für AWS Migration",
        "inhalt": [
            "Erstellung einer klaren Migrationsstrategie",
            "Durchführung von Proof-of-Concepts",
            "Automatisierung des Migrationsprozesses",
            "Kontinuierliche Überwachung und Optimierung",
            "Schulung und Vorbereitung des Teams"
        ]
    }
]

def move_slide(old_index, new_index):
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def add_slides(presentation, slides_data):
    # Wähle das Layout "Titel und Inhalt"
    slide_layout = None
    for layout in presentation.slide_layouts:
        if layout.name == "Titel und Inhalt weiß":
            slide_layout = layout
            break
    
    if not slide_layout:
        slide_layout = presentation.slide_layouts[1]
    
     # Füge die neuen Folien nach der Titelfolie ein
    for index, slide_data in enumerate(slides_data, start=1):
        # Füge eine neue Folie nach der Titelfolie ein
        new_slide = presentation.slides.add_slide(slide_layout)
        
        #verschiebe Folie vom Ende zur entsprechenden Position nach der Titelseite
        slide_id = presentation.slides.index(new_slide)
        move_slide(slide_id, index)

        
        # Setze den Titel
        title_shape = new_slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.get("titel", "")
        else:
            print(f"Warnung: Kein Titeltextfeld für Folie '{slide_data.get('titel', '')}'")
        
        # Füge den Inhalt hinzu
        content = new_slide.placeholders[1] if len(new_slide.placeholders) > 1 else None
        if content:
            tf = content.text_frame
            tf.clear()  # Entferne vorhandenen Text
            
            for punkt in slide_data.get("content_text", []):
                p = tf.add_paragraph()
                p.text = punkt
                p.level = 1

            
            # Formatiere den Text
            for paragraph in tf.paragraphs:
                if paragraph.font:
                    paragraph.font.size = Pt(18)
                    paragraph.space_before = Pt(6)
        else:
            print(f"Warnung: Kein Inhaltstextfeld für Folie '{slide_data.get('titel', '')}'")
    
    








#--Hauptbereich ---------------------------------------------------------------------------------------------------------------------------------------




#--Sidebar ---------------------------------------------------------------------------------------------------------------------------------------
st.sidebar.title("App-Steuerung")


#Schaltflächen für neues Dokument
st.sidebar.subheader("Neue Präsentation", divider='grey')
newdoc_form = st.sidebar.form("newdoc_form_key")
st.session_state.new_title = newdoc_form.text_input("Präsnetationstitel", value=st.session_state.new_title, help="Der Titel, den die Präsentation haben soll")

document_types = ["Sales Folien", "IT Konzept Folien"]
default_index = document_types.index(st.session_state.new_doctype) if st.session_state.new_doctype in document_types else 0
st.session_state.new_doctype = newdoc_form.selectbox("Folientyp", 
    options=document_types,
    index=default_index)

st.session_state.new_content_focus = newdoc_form.text_area("Inhaltlicher Schwerpunkt", value=st.session_state.new_content_focus, help="Nenne alle Aspekte, die in der Präsentation zwingend behandelt werden sollen.")
st.session_state.new_chapter_count = newdoc_form.slider("Anzahl der Kapitel.", min_value=1, max_value=30, value=st.session_state.new_chapter_count)

new_submitted = newdoc_form.form_submit_button("Foliensatz erstellen")

#Schaltflächen für den PPT Export
st.sidebar.subheader("PowerPoint Export", divider='grey')
if st.sidebar.button("PowerPoint Dokument generieren", key="ppt_export"):
    presentation = open_pptx_template()
    add_slides(presentation, st.session_state.toc_list)
    generate_ppt(presentation, st.session_state.new_title)



#--App Logik ---------------------------------------------------------------------------------------------------------------------------------------
# Foliensatz aus gegebenen Parametern per ChatBot erstellen lassen
if new_submitted:
    
    # Überschriften für Hauptbereich aus Parametern erzeugen
    st.session_state.new_header = st.session_state.new_doctype + ": " + st.session_state.new_title

    with st.spinner(text="Foliensatz wird erstellt ..."):
        # Inhaltsverzeichnis + Infotexte + Prompts aus Paramtetern per Chatbot erzeugen
        st.session_state.toc_list = generate_toc(st.session_state.new_doctype, st.session_state.new_title, st.session_state.new_content_focus, st.session_state.new_chapter_count)
    
    # Leere SessionState Elemente erzeugen die im Weiteren mit Inhalten gefüllt werden
    st.session_state.kapitel_header = [item["title"] for item in st.session_state.toc_list]
    st.session_state.kapitel_info = [item["help_text"] for item in st.session_state.toc_list]
    st.session_state.kapitel_inhalt = [item["content_text"] for item in st.session_state.toc_list]
    st.session_state.kapitel_prompt = [item["prompt_text"] for item in st.session_state.toc_list]
    st.session_state.image_prompt = [item["image_prompt_text"] for item in st.session_state.toc_list]


# Hauptbereich ------------------------------------------------
st.header(st.session_state.new_header, divider='grey') 

# Funktion zur Generierung von URL-freundlichen Ankern
def generate_anchor(text):
    # Entferne Umlaute und konvertiere zu ASCII
    text = unicodedata.normalize('NFKD', text).encode('ASCII', 'ignore').decode('ASCII')
    # Konvertiere zu Kleinbuchstaben und ersetze Leerzeichen durch Bindestriche
    text = re.sub(r'\W+', '-', text.lower()).strip('-')
    return text

# Funktion zur Erstellung des farbigen Inhaltsverzeichnisses
def create_colored_toc():
    st.subheader("Inhaltsverzeichnis")

    if not st.session_state.toc_list:
        st.info("""
        **Hinweis zur Erstellung des Inhaltsverzeichnisses:**
        
        Um ein Inhaltsverzeichnis zu erstellen, füllen Sie bitte die Felder im Bereich "Neue Präsentation" in der Seitenleiste aus und klicken Sie dann auf "Präsentation erstellen".
        
        Sobald die Folienstruktur generiert wurde, wird hier das Inhaltsverzeichnis angezeigt.
        """)
        return
    
    for i, item in enumerate(st.session_state.toc_list):
        title_text = item["title"]
        anchor = generate_anchor(title_text)
        
        # Überprüfen, ob Inhalt für dieses Kapitel vorhanden ist
        has_content = st.session_state.kapitel_inhalt[i].strip() != ""
        
        if has_content:
            color = "#d4edda"  # Grün (success)
            border_color = "#c3e6cb"
            icon = "✓"  # Grünes Häkchen
            icon_color = "#28a745"  # Grün
        else:
            color = "#cce5ff"  # Blau (info)
            border_color = "#b8daff"
            icon = "!"  # Ausrufezeichen
            icon_color = "#007bff"  # Blau
        
        # HTML für farbigen Link erstellen
        colored_link = f"""
        <div style="
            background-color: {color};
            border: 1px solid {border_color};
            border-radius: 5px;
            padding: 5px;
            margin: 2px 0;
            display: flex;
            justify-content: space-between;
            align-items: center;">
            <a href="#{anchor}" style="
                color: #333;
                text-decoration: none;
                font-weight: bold;">
                {title_text}
            </a>
            <span style="
                color: {icon_color};
                font-weight: bold;">
                {icon}
            </span>
        </div>
        """
        
        st.markdown(colored_link, unsafe_allow_html=True)

# Füge einen Anker für das Inhaltsverzeichnis hinzu
st.markdown('<a name="inhaltsverzeichnis"></a>', unsafe_allow_html=True)

# Inhaltsverzeichnis erstellen
create_colored_toc()

# Funktion zum Aktualisieren des Session States
def update_session_state(key, value):
    st.session_state[key] = value

# Erstellen der Webseiten-Struktur mit Überschriften, Infoboxen und Textboxen
for i, item in enumerate(st.session_state.toc_list):
    title_text = st.session_state.kapitel_header[i]
    anchor = generate_anchor(title_text)
    
    st.markdown(f'<a name="{anchor}"></a>', unsafe_allow_html=True)
    
    # Erstelle die Überschrift mit dem "Zurück zum Inhaltsverzeichnis" Icon
    st.markdown(f"""
    <h2 style="display: flex; justify-content: left; align-items: center;">
        {title_text}
        <a href="#inhaltsverzeichnis" style="text-decoration: none; color: inherit; font-size: 0.8em; margin-left: 5px;">
            &#128196; 
        </a>
    </h2>
    """, unsafe_allow_html=True)
    
    st.info(st.session_state.kapitel_info[i])
    st.session_state.kapitel_prompt[i] = st.text_area(f"Prompt zum generieren des Inhalts", value=st.session_state.kapitel_prompt[i], height=100)
    
    #if st.button("Kapitel " + title_text + " generieren", key=f"button_chapter_{i}"):
        #generate_chapter(title_text, st.session_state.kapitel_prompt[i], st.session_state.new_doctype, st.session_state.new_title, st.session_state.new_writing_style, st.session_state.new_word_count, st.session_state.new_context, st.session_state.new_stakeholder, i)
    
    st.session_state.kapitel_inhalt[i] = st.text_area(f"Textbaustein für {title_text}", value=st.session_state.kapitel_inhalt[i], height=300)





